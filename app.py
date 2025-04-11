from flask import Flask, request, jsonify, render_template, send_from_directory
import language_tool_python
import os
import pdfplumber
from docx import Document
import pandas as pd
from pptx import Presentation
from fpdf import FPDF  # Add this import for PDF creation

app = Flask(__name__)

UPLOAD_FOLDER = 'uploads'
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

# Initialize LanguageTool for English
tool = language_tool_python.LanguageTool('en-US')

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/check', methods=['POST'])
def check_text():
    data = request.form  # Change from request.json to request.form for FormData
    text = data.get('text', '')  # For text input
    file = request.files.get('file')  # For file uploads

    # If a file is uploaded, read its content
    if file:
        file_path = os.path.join(UPLOAD_FOLDER, file.filename)  # Define path
        file.save(file_path)  # Save the uploaded file
        text = extract_text_from_file(file_path)

    # Check for grammar and spelling errors
    matches = tool.check(text)
    errors = []
    highlighted_text = ""
    current_index = 0

    for match in matches:
        error_word = text[match.offset:match.offset + match.errorLength]
        errors.append({
            'word': error_word,
            'suggestion': ', '.join(match.replacements)
        })

        # Highlight text with error word in red
        highlighted_text += text[current_index:match.offset]
        highlighted_text += f'<span style="color: red;">{error_word}</span>'
        current_index = match.offset + match.errorLength

    highlighted_text += text[current_index:]  # Add remaining text

    # For file upload mode, create corrected PDF
    if file:
        pdf_path = os.path.join(UPLOAD_FOLDER, 'corrected_output.pdf')
        create_pdf(highlighted_text, pdf_path)
        return jsonify({'errors': errors, 'highlighted_text': highlighted_text, 'pdf_path': f'/download/{os.path.basename(pdf_path)}'})
    else:
        return jsonify({'errors': errors, 'highlighted_text': highlighted_text})

def create_pdf(text, pdf_path):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=10)
    pdf.set_font("Arial", size=12)
    pdf.multi_cell(0, 10, text.replace('<span style="color: red;">', '').replace('</span>', ''))
    pdf.output(pdf_path)

@app.route('/download/<filename>')
def download_file(filename):
    return send_from_directory(UPLOAD_FOLDER, filename)

def extract_text_from_file(file_path):
    text = ''
    if file_path.endswith('.txt'):
        with open(file_path, 'r', encoding='utf-8') as f:
            text = f.read()
    elif file_path.endswith('.docx'):
        doc = Document(file_path)
        for para in doc.paragraphs:
            text += para.text + '\n'
    elif file_path.endswith('.pdf'):
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text += page.extract_text() + '\n'
    elif file_path.endswith('.xlsx'):
        df = pd.read_excel(file_path)
        text = df.to_string(index=False)
    elif file_path.endswith('.pptx'):
        presentation = Presentation(file_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + '\n'
    else:
        return "Unsupported file type."
    os.remove(file_path)  # Remove the uploaded file after processing
    return text

if __name__ == '__main__':
    app.run(debug=True)
